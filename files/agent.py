"""
Data Analyst AI Agent (Free Gemini version)
=============================================
Ek AI-powered agent jo CSV/Excel/PDF/Word/PPT data ko clean, analyze, aur
query karta hai natural language (Hinglish/English) instructions se.
Uses Google Gemini's free tier (no billing needed).

Setup:
    pip install google-genai pandas openpyxl python-docx python-pptx pypdf

    Get a free key (no credit card): https://aistudio.google.com/apikey
    export GEMINI_API_KEY="your-key-here"

Usage:
    python agent.py your_data.csv
"""

import os
import sys
import io
import sqlite3
import traceback
import subprocess
import pandas as pd
from google import genai
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.utils import get_column_letter
from openpyxl.formatting.rule import CellIsRule, ColorScaleRule, FormulaRule
from openpyxl.worksheet.datavalidation import DataValidation
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

MODEL = "gemini-flash-latest"  # always points to current Gemini Flash model

client = genai.Client()  # reads GEMINI_API_KEY from environment


def call_llm(system_prompt: str, user_content: str, max_tokens: int = 800) -> str:
    """Wrapper so the rest of the code doesn't care which LLM provider is used."""
    response = client.models.generate_content(
        model=MODEL,
        contents=f"{system_prompt}\n\n{user_content}",
    )
    return response.text.strip()

# ----------------------------------------------------------------------
# 0. UNIVERSAL FILE READER (csv, xlsx, pdf, docx, pptx)
# ----------------------------------------------------------------------
def read_pdf_text(filepath: str) -> str:
    reader = PdfReader(filepath)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def read_docx_text(filepath: str) -> str:
    doc = DocxDocument(filepath)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def read_pptx_text(filepath: str) -> str:
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def load_any_file(filepath: str):
    """Returns ('table', DataFrame) for tabular files, or ('text', str) for documents."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".csv":
        return "table", pd.read_csv(filepath)
    elif ext in (".xlsx", ".xls"):
        return "table", pd.read_excel(filepath)
    elif ext == ".pdf":
        return "text", read_pdf_text(filepath)
    elif ext == ".docx":
        return "text", read_docx_text(filepath)
    elif ext == ".pptx":
        return "text", read_pptx_text(filepath)
    else:
        raise ValueError(f"File type '{ext}' abhi supported nahi hai.")


def ask_about_text(text: str, question: str) -> str:
    """NL Q&A over non-tabular documents (PDF/Word/PPT)."""
    truncated = text[:15000]  # keep prompt reasonable
    system_prompt = "Tum ek document analyst ho. User ke document content ke aadhar par uske sawaal ka Hinglish mein clear jawab do."
    user_content = f"Document content:\n{truncated}\n\nQuestion: {question}"
    return call_llm(system_prompt, user_content)


# ----------------------------------------------------------------------
# 1. DATA LOADING + AUTO CLEANING
# ----------------------------------------------------------------------
def load_data(filepath: str) -> pd.DataFrame:
    if filepath.endswith(".csv"):
        df = pd.read_csv(filepath)
    elif filepath.endswith((".xlsx", ".xls")):
        df = pd.read_excel(filepath)
    else:
        raise ValueError("Sirf CSV ya Excel files supported hain.")
    return df


def auto_clean(df: pd.DataFrame) -> tuple[pd.DataFrame, list[str]]:
    """Basic automated cleaning + log of what changed."""
    log = []
    before_rows = len(df)

    # Drop fully empty rows/cols
    df = df.dropna(how="all")
    df = df.dropna(axis=1, how="all")

    # Strip whitespace from string columns FIRST (so duplicate detection
    # and casing checks work on the real values, not on " x" vs "x")
    str_cols = df.select_dtypes(include="object").columns
    for col in str_cols:
        df[col] = df[col].astype(str).str.strip()
        df[col] = df[col].replace({"nan": pd.NA, "": pd.NA})

    # Remove exact duplicate rows (now whitespace-insensitive)
    dup_count = df.duplicated().sum()
    if dup_count > 0:
        df = df.drop_duplicates()
        log.append(f"{dup_count} duplicate rows hataye gaye")

    # Try converting object columns that look numeric/date
    for col in str_cols:
        converted = pd.to_numeric(df[col], errors="coerce")
        if converted.notna().sum() / max(len(df), 1) > 0.9:
            df[col] = converted

    after_rows = len(df)
    if after_rows != before_rows:
        log.append(f"Empty rows/cols clean kiye: {before_rows} -> {after_rows} rows")

    # Report missing values even if we didn't drop them (so the user isn't
    # told "sab clean hai" when there's real missing data sitting there)
    missing = df.isna().sum()
    missing = missing[missing > 0]
    if not missing.empty:
        details = ", ".join(f"{col}: {cnt}" for col, cnt in missing.items())
        log.append(f"Missing values mile (hataye nahi gaye, 'clean:' command se batao kya karna hai): {details}")

    # Flag columns with inconsistent casing (e.g. 'North' vs 'north')
    for col in str_cols:
        if col not in df.columns:
            continue
        non_null = df[col].dropna().astype(str)
        if non_null.empty:
            continue
        distinct = non_null.unique()
        distinct_lower = pd.Series(distinct).str.lower().unique()
        if len(distinct) > len(distinct_lower):
            log.append(f"'{col}' mein inconsistent casing mili (jaise 'North'/'north') — 'clean:' se normalize karwa sakte ho")

    if not log:
        log.append("Data pehle se hi clean tha, koi major issue nahi mila")

    return df, log


# ----------------------------------------------------------------------
# 2b. CUSTOM CLEANING (natural language instruction -> permanent df change)
# ----------------------------------------------------------------------
def custom_clean(df: pd.DataFrame, instruction: str):
    """Applies a user-described cleaning instruction and returns the UPDATED df."""
    schema = f"Columns: {list(df.columns)}\nDtypes:\n{df.dtypes.to_string()}\nSample rows:\n{df.head(3).to_string()}"

    system_prompt = (
        "Tum ek data cleaning AI ho. User ek pandas DataFrame `df` pe cleaning "
        "instruction dega (Hinglish/English). Tumhe SIRF valid Python pandas "
        "code dena hai jo `df` ko MODIFY kare aur naya df wapas `df` variable "
        "mein hi store kare (e.g. df = df[...] ya df['col'] = ...). "
        "Koi explanation, koi markdown fences, koi text nahi — sirf raw code."
    )

    user_content = f"DataFrame info:\n{schema}\n\nCleaning instruction: {instruction}\n\nCode:"
    code = call_llm(system_prompt, user_content, max_tokens=400)
    code = code.replace("```python", "").replace("```", "").strip()

    local_vars = {"df": df.copy(), "pd": pd}
    try:
        exec(code, {"__builtins__": {}}, local_vars)
        new_df = local_vars.get("df", df)
        return new_df, f"Cleaning applied:\n{code}\n\nNaya shape: {new_df.shape}"
    except Exception as e:
        return df, f"Error aaya cleaning code chalane mein:\n{code}\n\n{e}"


# ----------------------------------------------------------------------
# 2. NATURAL LANGUAGE -> PANDAS CODE (Q&A over data)
# ----------------------------------------------------------------------
def ask_question(df: pd.DataFrame, question: str) -> str:
    schema = f"Columns: {list(df.columns)}\nDtypes:\n{df.dtypes.to_string()}\nSample rows:\n{df.head(3).to_string()}"

    system_prompt = (
        "Tum ek data analyst AI ho. User ek pandas DataFrame `df` ke baare mein "
        "sawaal poochega (Hinglish ya English mein). Tumhe SIRF valid Python "
        "pandas code dena hai jo `result` variable mein answer store kare. "
        "Koi explanation, koi markdown fences, koi text nahi — sirf raw code."
    )

    user_content = f"DataFrame info:\n{schema}\n\nQuestion: {question}\n\nCode:"
    code = call_llm(system_prompt, user_content, max_tokens=500)
    code = code.replace("```python", "").replace("```", "").strip()

    # Sandbox execution with restricted scope
    local_vars = {"df": df, "pd": pd}
    try:
        exec(code, {"__builtins__": {}}, local_vars)
        result = local_vars.get("result", "Result variable nahi mila.")
        return f"Code executed:\n{code}\n\nAnswer:\n{result}"
    except Exception as e:
        return f"Error aaya code chalane mein:\n{code}\n\n{traceback.format_exc()}"


# ----------------------------------------------------------------------
# 3b. SAVE CLEANED DATA + PERSISTENT SQL DATABASE FILE
# ----------------------------------------------------------------------
def save_data(df: pd.DataFrame, base_path: str) -> str:
    """Saves current (possibly cleaned) df as both CSV and Excel."""
    csv_path = base_path + "_cleaned.csv"
    xlsx_path = base_path + "_cleaned.xlsx"
    df.to_csv(csv_path, index=False)
    df.to_excel(xlsx_path, index=False)
    return f"Cleaned data saved:\n  - {csv_path}\n  - {xlsx_path}"


def export_sql_db(df: pd.DataFrame, db_path: str) -> str:
    """Writes df into a REAL, persistent SQLite .db file that tools like
    SQL Workbench / DB Browser for SQLite can open directly."""
    conn = sqlite3.connect(db_path)
    df.to_sql("data", conn, index=False, if_exists="replace")
    conn.commit()
    conn.close()
    return (
        f"SQLite database saved: {db_path}\n"
        f"Table name: 'data'\n"
        f"Ise SQL Workbench / DB Browser for SQLite mein khol sakte ho — "
        f"CREATE TABLE aur data dono wahan dikhega."
    )


# ----------------------------------------------------------------------
# 3. SQL SUPPORT (load df into SQLite, run SQL queries)
# ----------------------------------------------------------------------
def run_sql(df: pd.DataFrame, sql_question: str) -> str:
    conn = sqlite3.connect(":memory:")
    df.to_sql("data", conn, index=False, if_exists="replace")

    schema = f"Table `data` with columns: {list(df.columns)}"
    system_prompt = (
        "Tum SQL expert ho. User ka sawaal SQLite query mein convert karo "
        "table `data` ke against. SIRF raw SQL do, koi explanation nahi."
    )

    user_content = f"{schema}\n\nQuestion: {sql_question}\n\nSQL:"
    sql = call_llm(system_prompt, user_content, max_tokens=300)
    sql = sql.replace("```sql", "").replace("```", "").strip()

    try:
        result_df = pd.read_sql_query(sql, conn)
        return f"SQL Query:\n{sql}\n\nResult:\n{result_df.to_string()}"
    except Exception as e:
        return f"SQL Error:\n{sql}\n\n{e}"
    finally:
        conn.close()


# ----------------------------------------------------------------------
# 4. AUTO EDA REPORT
# ----------------------------------------------------------------------
def generate_report(df: pd.DataFrame) -> str:
    summary = df.describe(include="all").to_string()
    missing = df.isna().sum()
    missing = missing[missing > 0].to_string() if missing.sum() > 0 else "Koi missing values nahi"

    system_prompt = (
        "Tum senior data analyst ho. Neeche diye gaye stats dekh kar 5-7 key "
        "insights Hinglish mein do, bullet points mein. Business-relevant "
        "patterns, anomalies, aur trends highlight karo."
    )

    user_content = f"Shape: {df.shape}\n\nSummary stats:\n{summary}\n\nMissing values:\n{missing}"
    return call_llm(system_prompt, user_content, max_tokens=600)


# ----------------------------------------------------------------------
# 5. EXCEL DASHBOARD (with real formulas + charts)
# ----------------------------------------------------------------------
def generate_excel_dashboard(df: pd.DataFrame, output_path: str) -> str:
    wb = Workbook()

    # --- Data sheet ---
    data_sheet = wb.active
    data_sheet.title = "Data"
    data_sheet.append(list(df.columns))
    for _, row in df.iterrows():
        data_sheet.append(list(row))

    header_fill = PatternFill("solid", start_color="4472C4")
    for cell in data_sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = header_fill
    for col_cells in data_sheet.columns:
        letter = get_column_letter(col_cells[0].column)
        data_sheet.column_dimensions[letter].width = 16

    n_rows = len(df) + 1  # +1 for header
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    cat_cols = df.select_dtypes(include="object").columns.tolist()

    # --- Summary sheet (formulas, not hardcoded values) ---
    summary = wb.create_sheet("Summary")
    summary["A1"] = "Metric"
    summary["B1"] = "Value"
    summary["A1"].font = Font(bold=True, name="Arial")
    summary["B1"].font = Font(bold=True, name="Arial")

    summary["A2"] = "Total Rows"
    summary["B2"] = f"=COUNTA(Data!A2:A{n_rows})"

    r = 3
    for col in numeric_cols:
        col_idx = df.columns.get_loc(col) + 1
        col_letter = get_column_letter(col_idx)
        rng = f"Data!{col_letter}2:{col_letter}{n_rows}"
        summary[f"A{r}"] = f"{col} - Sum"
        summary[f"B{r}"] = f"=SUM({rng})"
        summary[f"A{r+1}"] = f"{col} - Average"
        summary[f"B{r+1}"] = f"=AVERAGE({rng})"
        summary[f"A{r+2}"] = f"{col} - Max"
        summary[f"B{r+2}"] = f"=MAX({rng})"
        summary[f"A{r+3}"] = f"{col} - Min"
        summary[f"B{r+3}"] = f"=MIN({rng})"
        r += 4

    summary.column_dimensions["A"].width = 28
    summary.column_dimensions["B"].width = 16

    # --- Category breakdown sheet + bar chart ---
    if cat_cols:
        cat_col = cat_cols[0]
        counts = df[cat_col].value_counts()
        cat_sheet = wb.create_sheet("Breakdown")
        cat_sheet["A1"] = cat_col
        cat_sheet["B1"] = "Count"
        cat_sheet["A1"].font = Font(bold=True, name="Arial")
        cat_sheet["B1"].font = Font(bold=True, name="Arial")
        for i, (val, cnt) in enumerate(counts.items(), start=2):
            cat_sheet[f"A{i}"] = str(val)
            cat_sheet[f"B{i}"] = int(cnt)

        chart = BarChart()
        chart.title = f"{cat_col} Distribution"
        chart.x_axis.title = cat_col
        chart.y_axis.title = "Count"
        data_ref = Reference(cat_sheet, min_col=2, min_row=1, max_row=len(counts) + 1)
        cats_ref = Reference(cat_sheet, min_col=1, min_row=2, max_row=len(counts) + 1)
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)
        cat_sheet.add_chart(chart, "D2")

    # --- Trend line chart (if there's a numeric column to plot across rows) ---
    if numeric_cols:
        trend_col = numeric_cols[0]
        col_idx = df.columns.get_loc(trend_col) + 1
        line = LineChart()
        line.title = f"{trend_col} Trend"
        data_ref = Reference(data_sheet, min_col=col_idx, min_row=1, max_row=min(n_rows, 200))
        line.add_data(data_ref, titles_from_data=True)
        summary.add_chart(line, "D2")

    wb.save(output_path)

    # Recalculate formulas via LibreOffice
    try:
        subprocess.run(
            ["python", "/mnt/skills/public/xlsx/scripts/recalc.py", output_path, "30"],
            check=False, capture_output=True, timeout=40
        )
    except Exception:
        pass

    return f"Dashboard saved: {output_path}"


# ----------------------------------------------------------------------
# 8. ADVANCED EXCEL: FORMULAS, CONDITIONAL FORMATTING, DATA VALIDATION,
#    EDIT-EXISTING-FILE-IN-PLACE
# ----------------------------------------------------------------------
def add_formula_column(df: pd.DataFrame, instruction: str, output_path: str) -> str:
    """Adds a new column driven by an Excel formula (supports nested IF,
    INDEX-MATCH, VLOOKUP, etc.) — the formula is written as a real Excel
    formula, not a precomputed value, so it stays live in the file."""
    schema = f"Columns: {list(df.columns)}\nSample rows:\n{df.head(3).to_string()}"
    system_prompt = (
        "Tum Excel formula expert ho. User ek naya column chahta hai jiski "
        "value ek Excel formula se calculate ho (IF, nested IF, INDEX-MATCH, "
        "VLOOKUP, etc. use kar sakte ho). Reply STRICTLY is JSON format mein "
        '(sirf JSON, koi aur text nahi): {"column_name": "...", '
        '"formula_template": "=IF([Salary]>100000,\\"Yes\\",\\"No\\")"} '
        "formula_template mein column names ko [ColumnName] format mein likho "
        "— main unhe actual Excel cell references (jaise B2) se replace kar dunga."
    )
    raw = call_llm(system_prompt, f"{schema}\n\nInstruction: {instruction}", max_tokens=400)
    raw = raw.replace("```json", "").replace("```", "").strip()

    import json
    try:
        spec = json.loads(raw)
    except Exception as e:
        return f"Formula spec samajh nahi aayi:\n{raw}\n\nError: {e}"

    col_name = spec["column_name"]
    template = spec["formula_template"]

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    headers = list(df.columns) + [col_name]
    ws.append(headers)
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")

    col_letter = {col: get_column_letter(i + 1) for i, col in enumerate(df.columns)}

    for r, (_, row) in enumerate(df.iterrows(), start=2):
        row_values = list(row)
        formula = template
        for col, letter in col_letter.items():
            formula = formula.replace(f"[{col}]", f"{letter}{r}")
        row_values.append(formula)
        ws.append(row_values)

    wb.save(output_path)
    try:
        subprocess.run(
            ["python", "/mnt/skills/public/xlsx/scripts/recalc.py", output_path, "30"],
            check=False, capture_output=True, timeout=40
        )
    except Exception:
        pass

    return f"Formula column '{col_name}' added, file saved: {output_path}\nFormula used: {template}"


def add_conditional_formatting(df: pd.DataFrame, instruction: str, output_path: str) -> str:
    """Applies conditional formatting (highlight cells, color scale) based
    on a natural-language rule, e.g. 'highlight Sales below 0 in red'."""
    schema = f"Columns: {list(df.columns)}"
    system_prompt = (
        "Tum Excel conditional-formatting expert ho. Reply STRICTLY JSON mein "
        '(sirf JSON): {"column": "Sales", "rule_type": "cell_is", '
        '"operator": "lessThan", "value": 0, "color": "FFC7CE"} '
        'rule_type "cell_is" (operator: lessThan/greaterThan/equal, value chahiye) '
        'ya "color_scale" (koi value nahi chahiye) ho sakta hai. color ek 6-char '
        "hex bina # ke."
    )
    raw = call_llm(system_prompt, f"{schema}\n\nInstruction: {instruction}", max_tokens=300)
    raw = raw.replace("```json", "").replace("```", "").strip()

    import json
    try:
        spec = json.loads(raw)
    except Exception as e:
        return f"Rule samajh nahi aayi:\n{raw}\n\nError: {e}"

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(list(df.columns))
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")
    for _, row in df.iterrows():
        ws.append(list(row))

    col = spec["column"]
    col_idx = df.columns.get_loc(col) + 1
    col_letter = get_column_letter(col_idx)
    n_rows = len(df) + 1
    cell_range = f"{col_letter}2:{col_letter}{n_rows}"

    if spec["rule_type"] == "color_scale":
        rule = ColorScaleRule(
            start_type="min", start_color="F8696B",
            end_type="max", end_color="63BE7B"
        )
        ws.conditional_formatting.add(cell_range, rule)
    else:
        fill = PatternFill(start_color=spec["color"], end_color=spec["color"], fill_type="solid")
        rule = CellIsRule(operator=spec["operator"], formula=[str(spec["value"])], fill=fill)
        ws.conditional_formatting.add(cell_range, rule)

    wb.save(output_path)
    return f"Conditional formatting applied on '{col}', file saved: {output_path}"


def add_data_validation(df: pd.DataFrame, instruction: str, output_path: str) -> str:
    """Adds a dropdown (data validation) to a column, e.g. restricting
    Status to Active/Inactive/Pending."""
    schema = f"Columns: {list(df.columns)}\nSample:\n{df.head(3).to_string()}"
    system_prompt = (
        "Tum Excel data-validation expert ho. Reply STRICTLY JSON mein "
        '(sirf JSON): {"column": "Status", "options": ["Active","Inactive","Pending"]}'
    )
    raw = call_llm(system_prompt, f"{schema}\n\nInstruction: {instruction}", max_tokens=300)
    raw = raw.replace("```json", "").replace("```", "").strip()

    import json
    try:
        spec = json.loads(raw)
    except Exception as e:
        return f"Validation spec samajh nahi aayi:\n{raw}\n\nError: {e}"

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(list(df.columns))
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")
    for _, row in df.iterrows():
        ws.append(list(row))

    col = spec["column"]
    col_idx = df.columns.get_loc(col) + 1
    col_letter = get_column_letter(col_idx)
    n_rows = len(df) + 1
    options_str = ",".join(spec["options"])

    dv = DataValidation(type="list", formula1=f'"{options_str}"', allow_blank=True)
    ws.add_data_validation(dv)
    dv.add(f"{col_letter}2:{col_letter}{n_rows}")

    wb.save(output_path)
    return f"Dropdown validation added on '{col}' ({', '.join(spec['options'])}), file saved: {output_path}"


def edit_excel_in_place(source_path: str, df: pd.DataFrame, output_path: str) -> str:
    """Loads the ORIGINAL Excel file (preserving its sheets, styles, other
    tabs) and updates the main data sheet with the current (possibly
    cleaned) DataFrame, then saves as a new file so the original is safe."""
    if not source_path.lower().endswith((".xlsx", ".xls")):
        return "Ye command sirf tab kaam karta hai jab original file Excel (.xlsx) ho."

    wb = load_workbook(source_path)
    ws = wb.active  # first/main sheet

    # Clear existing data rows (keep other sheets untouched)
    ws.delete_rows(1, ws.max_row)

    ws.append(list(df.columns))
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")
    for _, row in df.iterrows():
        ws.append(list(row))

    wb.save(output_path)
    sheet_names = wb.sheetnames
    return (
        f"Excel file updated (saved as new file to keep original safe): {output_path}\n"
        f"Sheets preserved: {sheet_names}"
    )


# ----------------------------------------------------------------------
# 6. POWER BI EXPORT (clean, typed data ready for import)
# ----------------------------------------------------------------------
def export_for_powerbi(df: pd.DataFrame, output_path: str) -> str:
    export_df = df.copy()
    for col in export_df.select_dtypes(include="object").columns:
        parsed = pd.to_datetime(export_df[col], errors="coerce")
        if parsed.notna().sum() / max(len(export_df), 1) > 0.9:
            export_df[col] = parsed

    export_df.columns = [str(c).strip().replace(" ", "_") for c in export_df.columns]
    export_df.to_excel(output_path, index=False, sheet_name="Data")
    return f"Power BI-ready file saved: {output_path} (clean headers, typed columns, single flat table)"


# ----------------------------------------------------------------------
# 7. POWER BI: DAX MEASURES + DASHBOARD SETUP GUIDE
# ----------------------------------------------------------------------
def generate_powerbi_guide(df: pd.DataFrame, output_path: str) -> str:
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    cat_cols = df.select_dtypes(include="object").columns.tolist()
    schema = f"Columns: {list(df.columns)}\nNumeric: {numeric_cols}\nCategorical: {cat_cols}\nSample:\n{df.head(3).to_string()}"

    system_prompt = (
        "Tum Power BI expert consultant ho. Data schema dekh kar do:\n"
        "1. 5-8 useful DAX measures (copy-paste ready code blocks, har ek ka short explanation)\n"
        "2. Step-by-step guide (Hinglish) ki Power BI Desktop mein is data se dashboard kaise banaye — "
        "konse visuals (card, bar, line, slicer, matrix) kahan lagayein, konsa column kahan drag karein\n"
        "3. Layout suggestion (top row KPIs, middle charts, bottom filters jaisa)\n"
        "Markdown format mein do, clear headers ke saath."
    )

    guide = call_llm(system_prompt, f"Data schema:\n{schema}", max_tokens=1500)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(guide)
    return f"Power BI guide saved: {output_path}\n\n{guide}"


# ----------------------------------------------------------------------
# MAIN CLI LOOP
# ----------------------------------------------------------------------
def main():
    if len(sys.argv) < 2:
        print("Usage: python agent.py <file.csv|xlsx|pdf|docx|pptx>")
        return

    filepath = sys.argv[1]
    print(f"Loading {filepath}...")
    kind, content = load_any_file(filepath)

    if kind == "text":
        print(f"Document loaded ({len(content)} characters)")
        print("\n--- Data Analyst AI Agent ready (document mode) ---")
        print("Commands: '<any question>' | 'exit'\n")
        while True:
            user_input = input("You: ").strip()
            if user_input.lower() in ("exit", "quit"):
                break
            print("\n" + ask_about_text(content, user_input) + "\n")
        return

    # kind == "table"
    df = content
    print(f"Loaded: {df.shape[0]} rows, {df.shape[1]} columns")

    df, log = auto_clean(df)
    print("\nCleaning log:")
    for entry in log:
        print(f"  - {entry}")

    print("\n--- Data Analyst AI Agent ready ---")
    print("Commands: 'report' | 'sql: <question>' | 'clean: <instruction>' | 'save' | 'sql-export' | 'dashboard' | 'powerbi' | 'powerbi-guide' | 'formula: <instruction>' | 'format: <instruction>' | 'validate: <instruction>' | 'edit-excel' | '<any question>' | 'exit'\n")

    while True:
        user_input = input("You: ").strip()
        if user_input.lower() in ("exit", "quit"):
            break
        elif user_input.lower() == "report":
            print("\n" + generate_report(df) + "\n")
        elif user_input.lower().startswith("sql:"):
            print("\n" + run_sql(df, user_input[4:].strip()) + "\n")
        elif user_input.lower().startswith("clean:"):
            df, msg = custom_clean(df, user_input[6:].strip())
            print("\n" + msg + "\n")
        elif user_input.lower() == "save":
            base = os.path.splitext(filepath)[0]
            print("\n" + save_data(df, base) + "\n")
        elif user_input.lower() == "sql-export":
            db_path = os.path.splitext(filepath)[0] + ".db"
            print("\n" + export_sql_db(df, db_path) + "\n")
        elif user_input.lower() == "dashboard":
            out = os.path.splitext(filepath)[0] + "_dashboard.xlsx"
            print("\n" + generate_excel_dashboard(df, out) + "\n")
        elif user_input.lower() == "powerbi":
            out = os.path.splitext(filepath)[0] + "_powerbi.xlsx"
            print("\n" + export_for_powerbi(df, out) + "\n")
        elif user_input.lower() == "powerbi-guide":
            out = os.path.splitext(filepath)[0] + "_powerbi_guide.md"
            print("\n" + generate_powerbi_guide(df, out) + "\n")
        elif user_input.lower().startswith("formula:"):
            out = os.path.splitext(filepath)[0] + "_formula.xlsx"
            print("\n" + add_formula_column(df, user_input[8:].strip(), out) + "\n")
        elif user_input.lower().startswith("format:"):
            out = os.path.splitext(filepath)[0] + "_formatted.xlsx"
            print("\n" + add_conditional_formatting(df, user_input[7:].strip(), out) + "\n")
        elif user_input.lower().startswith("validate:"):
            out = os.path.splitext(filepath)[0] + "_validated.xlsx"
            print("\n" + add_data_validation(df, user_input[9:].strip(), out) + "\n")
        elif user_input.lower() == "edit-excel":
            out = os.path.splitext(filepath)[0] + "_edited.xlsx"
            print("\n" + edit_excel_in_place(filepath, df, out) + "\n")
        else:
            print("\n" + ask_question(df, user_input) + "\n")


if __name__ == "__main__":
    main()