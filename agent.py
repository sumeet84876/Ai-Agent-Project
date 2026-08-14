"""
Data Analyst AI Agent (Free Gemini version)
=============================================
An AI-powered agent that cleans, analyzes, and queries CSV/Excel/PDF/Word/PPT
data using natural language instructions.
Uses Google Gemini's free tier (no billing needed).

Setup:
    pip install google-genai pandas openpyxl python-docx python-pptx pypdf

    Get a free key (no credit card): https://aistudio.google.com/apikey
    export GEMINI_API_KEY="your-key-here"

Usage:
    python agent.py your_data.csv
"""

import os
import sys
import io
import sqlite3
import traceback
import pandas as pd
from google import genai
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.utils import get_column_letter
from openpyxl.formatting.rule import CellIsRule, ColorScaleRule, FormulaRule
from openpyxl.worksheet.datavalidation import DataValidation
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

MODEL = "gemini-flash-latest"  # always points to current Gemini Flash model

# Language rule: ONLY free-form chat replies (general questions, document
# Q&A) are in Hinglish. Everything else — reports, guides, formulas, SQL,
# error messages, code — stays in English, since those are meant to be
# read/reused/shown professionally, not just chatted about.
HINGLISH_INSTRUCTION = " Respond in Hinglish (a natural mix of Hindi and English, written in Latin/Roman script)."
NO_FILLER_INSTRUCTION = " Be direct and concise. Do not use filler phrases like 'Great question!', 'Sure, I'd be happy to help', or unnecessary pleasantries — just answer."

# Safe builtins for the exec() sandbox used by custom_clean/ask_question.
# An EMPTY builtins dict (the original approach) blocks even harmless,
# extremely common calls like len(), sum(), round() that Gemini generates
# constantly in normal pandas code — causing frequent, confusing failures.
# This allowlist keeps dangerous stuff (open, __import__, eval, exec, input,
# exit) out while letting normal data-manipulation code actually run.
_SAFE_BUILTIN_NAMES = [
    "len", "sum", "min", "max", "round", "abs", "sorted", "reversed",
    "enumerate", "range", "zip", "map", "filter", "list", "dict", "set",
    "tuple", "str", "int", "float", "bool", "isinstance", "type",
    "True", "False", "None", "print",
]
SAFE_BUILTINS = {name: getattr(__builtins__, name, None) if not isinstance(__builtins__, dict)
                  else __builtins__.get(name) for name in _SAFE_BUILTIN_NAMES}
SAFE_BUILTINS = {k: v for k, v in SAFE_BUILTINS.items() if v is not None}
SAFE_BUILTINS.update({"True": True, "False": False, "None": None})

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOADS_DIR = os.path.join(BASE_DIR, "uploads")
OUTPUTS_DIR = os.path.join(BASE_DIR, "outputs")
os.makedirs(UPLOADS_DIR, exist_ok=True)
os.makedirs(OUTPUTS_DIR, exist_ok=True)

client = genai.Client()  # reads GEMINI_API_KEY from environment


def call_llm(system_prompt: str, user_content: str, max_tokens: int = 800) -> str:
    """Wrapper so the rest of the code doesn't care which LLM provider is used."""
    response = client.models.generate_content(
        model=MODEL,
        contents=f"{system_prompt}\n\n{user_content}",
    )
    return response.text.strip()


def generate_and_execute_with_retry(system_prompt: str, user_content: str, executor_fn,
                                     max_attempts: int = 3, max_tokens: int = 500):
    """Generic self-correcting agent loop: Analyze (context passed in) ->
    Plan (LLM generates code) -> Execute (executor_fn runs it) ->
    Evaluate (did it succeed?) -> Adapt (feed the error back and regenerate)
    -> repeat until success or max_attempts reached.

    executor_fn(code) must return (success: bool, result_or_error_message).
    Returns (final_code, result_or_error, success: bool, attempts_used: int).
    """
    code = None
    last_error = None

    for attempt in range(1, max_attempts + 1):
        if attempt == 1:
            code = call_llm(system_prompt, user_content, max_tokens=max_tokens)
        else:
            retry_content = (
                f"{user_content}\n\n"
                f"--- Previous attempt (#{attempt - 1}) ---\n{code}\n\n"
                f"--- That attempt failed with this error ---\n{last_error}\n\n"
                f"Fix the code so it works correctly. Return ONLY the corrected code, "
                f"nothing else."
            )
            code = call_llm(system_prompt, retry_content, max_tokens=max_tokens)

        code = code.replace("```python", "").replace("```", "").replace("```sql", "").strip()

        success, result = executor_fn(code)
        if success:
            return code, result, True, attempt
        last_error = result

    return code, last_error, False, max_attempts

# ----------------------------------------------------------------------
# 0. UNIVERSAL FILE READER (csv, xlsx, pdf, docx, pptx)
# ----------------------------------------------------------------------
def read_pdf_text(filepath: str) -> str:
    reader = PdfReader(filepath)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def read_docx_text(filepath: str) -> str:
    doc = DocxDocument(filepath)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def read_pptx_text(filepath: str) -> str:
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def load_any_file(filepath: str):
    """Returns ('table', DataFrame) for tabular files, or ('text', str) for documents."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".csv":
        return "table", pd.read_csv(filepath)
    elif ext in (".xlsx", ".xls"):
        return "table", pd.read_excel(filepath)
    elif ext == ".pdf":
        return "text", read_pdf_text(filepath)
    elif ext == ".docx":
        return "text", read_docx_text(filepath)
    elif ext == ".pptx":
        return "text", read_pptx_text(filepath)
    else:
        raise ValueError(f"File type '{ext}' is not supported yet.")


def ask_about_text(text: str, question: str) -> str:
    """NL Q&A over non-tabular documents (PDF/Word/PPT)."""
    truncated = text[:15000]  # keep prompt reasonable
    system_prompt = "You are a document analyst AI. Answer the user's question clearly, based on the document content provided." + HINGLISH_INSTRUCTION + NO_FILLER_INSTRUCTION
    user_content = f"Document content:\n{truncated}\n\nQuestion: {question}"
    return call_llm(system_prompt, user_content)


# ----------------------------------------------------------------------
# 1. DATA LOADING + AUTO CLEANING
# ----------------------------------------------------------------------
def load_data(filepath: str) -> pd.DataFrame:
    if filepath.endswith(".csv"):
        df = pd.read_csv(filepath)
    elif filepath.endswith((".xlsx", ".xls")):
        df = pd.read_excel(filepath)
    else:
        raise ValueError("Only CSV or Excel files are supported.")
    return df


def auto_clean(df: pd.DataFrame) -> tuple[pd.DataFrame, list[str]]:
    """Basic automated cleaning + log of what changed."""
    log = []
    before_rows = len(df)

    # Drop fully empty rows/cols
    df = df.dropna(how="all")
    df = df.dropna(axis=1, how="all")

    # Strip whitespace from string columns FIRST (so duplicate detection
    # and casing checks work on the real values, not on " x" vs "x")
    str_cols = df.select_dtypes(include="object").columns
    for col in str_cols:
        df[col] = df[col].astype(str).str.strip()
        df[col] = df[col].replace({"nan": pd.NA, "": pd.NA})

    # Remove exact duplicate rows (now whitespace-insensitive)
    dup_count = df.duplicated().sum()
    if dup_count > 0:
        df = df.drop_duplicates()
        log.append(f"Removed {dup_count} duplicate rows")

    # Try converting object columns that look numeric/date
    for col in str_cols:
        converted = pd.to_numeric(df[col], errors="coerce")
        if converted.notna().sum() / max(len(df), 1) > 0.9:
            df[col] = converted

    after_rows = len(df)
    if after_rows != before_rows:
        log.append(f"Cleaned empty rows/columns: {before_rows} -> {after_rows} rows")

    # Report missing values even if we didn't drop them (so the user isn't
    # told "everything's clean" when there's real missing data sitting there)
    missing = df.isna().sum()
    missing = missing[missing > 0]
    if not missing.empty:
        details = ", ".join(f"{col}: {cnt}" for col, cnt in missing.items())
        log.append(f"Missing values found (not removed, use 'clean:' to specify what to do): {details}")

    # Flag columns with inconsistent casing (e.g. 'North' vs 'north')
    for col in str_cols:
        if col not in df.columns:
            continue
        non_null = df[col].dropna().astype(str)
        if non_null.empty:
            continue
        distinct = non_null.unique()
        distinct_lower = pd.Series(distinct).str.lower().unique()
        if len(distinct) > len(distinct_lower):
            log.append(f"'{col}' has inconsistent casing (e.g. 'North'/'north') — use 'clean:' to normalize")

    if not log:
        log.append("Data was already clean, no major issues found")

    return df, log


# ----------------------------------------------------------------------
# 2b. CUSTOM CLEANING (natural language instruction -> permanent df change)
# ----------------------------------------------------------------------
def custom_clean(df: pd.DataFrame, instruction: str):
    """Applies a user-described cleaning instruction and returns the UPDATED df.
    Uses a self-correcting loop: if the generated code errors out, the error
    is fed back to the model so it can fix and retry (up to 3 attempts)."""
    schema = f"Columns: {list(df.columns)}\nDtypes:\n{df.dtypes.to_string()}\nSample rows:\n{df.head(3).to_string()}"

    system_prompt = (
        "You are a data cleaning AI. The user will give a cleaning "
        "instruction for a pandas DataFrame `df`. Respond with ONLY valid "
        "Python pandas code that MODIFIES `df` and stores the new result "
        "back in the `df` variable (e.g. df = df[...] or df['col'] = ...). "
        "No explanation, no markdown fences, no text — just raw code."
    )
    user_content = f"DataFrame info:\n{schema}\n\nCleaning instruction: {instruction}\n\nCode:"

    def executor(code: str):
        local_vars = {"df": df.copy(), "pd": pd}
        try:
            exec(code, {"__builtins__": SAFE_BUILTINS}, local_vars)
            new_df = local_vars.get("df", df)
            return True, new_df
        except Exception as e:
            return False, str(e)

    code, result, success, attempts = generate_and_execute_with_retry(
        system_prompt, user_content, executor, max_attempts=3, max_tokens=400
    )

    if success:
        note = f" (fixed after {attempts} attempts)" if attempts > 1 else ""
        return result, f"Cleaning applied{note}:\n{code}\n\nNew shape: {result.shape}"
    else:
        return df, f"Could not complete cleaning after {attempts} attempts.\nLast code:\n{code}\n\nLast error:\n{result}"


# ----------------------------------------------------------------------
# 2. NATURAL LANGUAGE -> PANDAS CODE (Q&A over data)
# ----------------------------------------------------------------------
def ask_question(df: pd.DataFrame, question: str) -> str:
    """Answers a natural-language question about the data using a
    self-correcting loop: if the generated code errors (or produces no
    result), the error is fed back so the model can fix and retry."""
    schema = f"Columns: {list(df.columns)}\nDtypes:\n{df.dtypes.to_string()}\nSample rows:\n{df.head(3).to_string()}"

    system_prompt = (
        "You are a data analyst AI. The user will ask a question about a "
        "pandas DataFrame `df`. Respond with ONLY valid Python pandas code "
        "that stores the answer in a `result` variable. No explanation, no "
        "markdown fences, no text — just raw code."
    )
    user_content = f"DataFrame info:\n{schema}\n\nQuestion: {question}\n\nCode:"

    def executor(code: str):
        local_vars = {"df": df.copy(), "pd": pd}
        try:
            exec(code, {"__builtins__": SAFE_BUILTINS}, local_vars)
            if "result" not in local_vars:
                return False, "Code ran but did not set a 'result' variable."
            return True, local_vars["result"]
        except Exception:
            return False, traceback.format_exc()

    code, result, success, attempts = generate_and_execute_with_retry(
        system_prompt, user_content, executor, max_attempts=3, max_tokens=500
    )

    if success:
        note = f" (fixed after {attempts} attempts)" if attempts > 1 else ""
        return f"Code executed{note}:\n{code}\n\nAnswer:\n{result}"
    else:
        return f"Could not answer after {attempts} attempts.\nLast code:\n{code}\n\nLast error:\n{result}"


# ----------------------------------------------------------------------
# 3b. SAVE CLEANED DATA + PERSISTENT SQL DATABASE FILE
# ----------------------------------------------------------------------
def save_data(df: pd.DataFrame, base_path: str) -> str:
    """Saves current (possibly cleaned) df as both CSV and Excel."""
    csv_path = base_path + "_cleaned.csv"
    xlsx_path = base_path + "_cleaned.xlsx"
    df.to_csv(csv_path, index=False)
    df.to_excel(xlsx_path, index=False)
    return f"Cleaned data saved:\n  - {csv_path}\n  - {xlsx_path}"


def export_sql_db(df: pd.DataFrame, db_path: str) -> str:
    """Writes df into a REAL, persistent SQLite .db file that tools like
    SQL Workbench / DB Browser for SQLite can open directly."""
    conn = sqlite3.connect(db_path)
    df.to_sql("data", conn, index=False, if_exists="replace")
    conn.commit()
    conn.close()
    return (
        f"SQLite database saved: {db_path}\n"
        f"Table name: 'data'\n"
        f"You can open this in SQL Workbench / DB Browser for SQLite — "
        f"both the CREATE TABLE statement and the data will be visible there."
    )


# ----------------------------------------------------------------------
# 3. SQL SUPPORT (load df into SQLite, run SQL queries)
# ----------------------------------------------------------------------
def run_sql(df: pd.DataFrame, sql_question: str) -> str:
    """Converts a natural-language question into SQL and runs it, using a
    self-correcting loop: if the query errors out, the error is fed back
    so the model can fix and retry."""
    conn = sqlite3.connect(":memory:")
    df.to_sql("data", conn, index=False, if_exists="replace")

    schema = f"Table `data` with columns: {list(df.columns)}"
    system_prompt = (
        "You are a SQL expert. Convert the user's question into a SQLite "
        "query against the table `data`. Respond with ONLY the raw SQL, "
        "no explanation."
    )
    user_content = f"{schema}\n\nQuestion: {sql_question}\n\nSQL:"

    def executor(sql: str):
        try:
            result_df = pd.read_sql_query(sql, conn)
            return True, result_df
        except Exception as e:
            return False, str(e)

    sql, result, success, attempts = generate_and_execute_with_retry(
        system_prompt, user_content, executor, max_attempts=3, max_tokens=300
    )
    conn.close()

    if success:
        note = f" (fixed after {attempts} attempts)" if attempts > 1 else ""
        return f"SQL Query{note}:\n{sql}\n\nResult:\n{result.to_string()}"
    else:
        return f"Could not run query after {attempts} attempts.\nLast SQL:\n{sql}\n\nLast error:\n{result}"


# ----------------------------------------------------------------------
# 4. AUTO EDA REPORT
# ----------------------------------------------------------------------
def generate_report(df: pd.DataFrame) -> str:
    summary = df.describe(include="all").to_string()
    missing = df.isna().sum()
    missing = missing[missing > 0].to_string() if missing.sum() > 0 else "No missing values"

    system_prompt = (
        "You are a senior data analyst. Based on the stats below, provide "
        "5-7 key insights as bullet points, plus a short 'Issues found' "
        "section listing data quality problems (missing values, outliers, "
        "inconsistent types) and a 'Suggested changes' section with concrete "
        "fixes. Highlight business-relevant patterns, anomalies, and trends. "
        "Respond in English — this report is a professional document, not "
        "a chat reply." + NO_FILLER_INSTRUCTION
    )

    user_content = f"Shape: {df.shape}\n\nSummary stats:\n{summary}\n\nMissing values:\n{missing}"
    return call_llm(system_prompt, user_content, max_tokens=800)


# ----------------------------------------------------------------------
# 5. EXCEL DASHBOARD (with real formulas + charts)
# ----------------------------------------------------------------------
def generate_excel_dashboard(df: pd.DataFrame, output_path: str) -> str:
    wb = Workbook()

    # --- Data sheet ---
    data_sheet = wb.active
    data_sheet.title = "Data"
    data_sheet.append(list(df.columns))
    for _, row in df.iterrows():
        data_sheet.append(list(row))

    header_fill = PatternFill("solid", start_color="4472C4")
    for cell in data_sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = header_fill
    for col_cells in data_sheet.columns:
        letter = get_column_letter(col_cells[0].column)
        data_sheet.column_dimensions[letter].width = 16

    n_rows = len(df) + 1  # +1 for header
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    cat_cols = df.select_dtypes(include="object").columns.tolist()

    # --- Summary sheet (formulas, not hardcoded values) ---
    summary = wb.create_sheet("Summary")
    summary["A1"] = "Metric"
    summary["B1"] = "Value"
    summary["A1"].font = Font(bold=True, name="Arial")
    summary["B1"].font = Font(bold=True, name="Arial")

    summary["A2"] = "Total Rows"
    summary["B2"] = f"=COUNTA(Data!A2:A{n_rows})"

    r = 3
    for col in numeric_cols:
        col_idx = df.columns.get_loc(col) + 1
        col_letter = get_column_letter(col_idx)
        rng = f"Data!{col_letter}2:{col_letter}{n_rows}"
        summary[f"A{r}"] = f"{col} - Sum"
        summary[f"B{r}"] = f"=SUM({rng})"
        summary[f"A{r+1}"] = f"{col} - Average"
        summary[f"B{r+1}"] = f"=AVERAGE({rng})"
        summary[f"A{r+2}"] = f"{col} - Max"
        summary[f"B{r+2}"] = f"=MAX({rng})"
        summary[f"A{r+3}"] = f"{col} - Min"
        summary[f"B{r+3}"] = f"=MIN({rng})"
        r += 4

    summary.column_dimensions["A"].width = 28
    summary.column_dimensions["B"].width = 16

    # --- Category breakdown sheet + bar chart ---
    if cat_cols:
        cat_col = cat_cols[0]
        counts = df[cat_col].value_counts()
        cat_sheet = wb.create_sheet("Breakdown")
        cat_sheet["A1"] = cat_col
        cat_sheet["B1"] = "Count"
        cat_sheet["A1"].font = Font(bold=True, name="Arial")
        cat_sheet["B1"].font = Font(bold=True, name="Arial")
        for i, (val, cnt) in enumerate(counts.items(), start=2):
            cat_sheet[f"A{i}"] = str(val)
            cat_sheet[f"B{i}"] = int(cnt)

        chart = BarChart()
        chart.title = f"{cat_col} Distribution"
        chart.x_axis.title = cat_col
        chart.y_axis.title = "Count"
        data_ref = Reference(cat_sheet, min_col=2, min_row=1, max_row=len(counts) + 1)
        cats_ref = Reference(cat_sheet, min_col=1, min_row=2, max_row=len(counts) + 1)
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)
        cat_sheet.add_chart(chart, "D2")

    # --- Trend line chart (if there's a numeric column to plot across rows) ---
    if numeric_cols:
        trend_col = numeric_cols[0]
        col_idx = df.columns.get_loc(trend_col) + 1
        line = LineChart()
        line.title = f"{trend_col} Trend"
        data_ref = Reference(data_sheet, min_col=col_idx, min_row=1, max_row=min(n_rows, 200))
        line.add_data(data_ref, titles_from_data=True)
        summary.add_chart(line, "D2")

    wb.save(output_path)
    # Note: formulas are written live (not pre-computed) — Excel/LibreOffice
    # recalculates them automatically the first time the file is opened.
    return f"Dashboard saved: {output_path}"


# ----------------------------------------------------------------------
# 8. ADVANCED EXCEL: FORMULAS, CONDITIONAL FORMATTING, DATA VALIDATION,
#    EDIT-EXISTING-FILE-IN-PLACE
# ----------------------------------------------------------------------
def add_formula_column(df: pd.DataFrame, instruction: str, output_path: str) -> str:
    """Adds a new column driven by an Excel formula (supports nested IF,
    INDEX-MATCH, VLOOKUP, etc.) — the formula is written as a real Excel
    formula, not a precomputed value, so it stays live in the file."""
    schema = f"Columns: {list(df.columns)}\nSample rows:\n{df.head(3).to_string()}"
    system_prompt = (
        "You are an Excel formula expert. The user wants a new column whose "
        "value is calculated by an Excel formula (you can use IF, nested IF, "
        "INDEX-MATCH, VLOOKUP, etc). Reply STRICTLY in this JSON format "
        '(JSON only, no other text): {"column_name": "...", '
        '"formula_template": "=IF([Salary]>100000,\\"Yes\\",\\"No\\")"} '
        "In formula_template, write column names as [ColumnName] — "
        "I will replace them with actual Excel cell references (e.g. B2)."
    )
    raw = call_llm(system_prompt, f"{schema}\n\nInstruction: {instruction}", max_tokens=400)
    raw = raw.replace("```json", "").replace("```", "").strip()

    import json
    try:
        spec = json.loads(raw)
    except Exception as e:
        return f"Could not parse formula spec:\n{raw}\n\nError: {e}"

    col_name = spec["column_name"]
    template = spec["formula_template"]

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    headers = list(df.columns) + [col_name]
    ws.append(headers)
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")

    col_letter = {col: get_column_letter(i + 1) for i, col in enumerate(df.columns)}

    for r, (_, row) in enumerate(df.iterrows(), start=2):
        row_values = list(row)
        formula = template
        for col, letter in col_letter.items():
            formula = formula.replace(f"[{col}]", f"{letter}{r}")
        row_values.append(formula)
        ws.append(row_values)

    wb.save(output_path)
    return f"Formula column '{col_name}' added, file saved: {output_path}\nFormula used: {template}"


def add_conditional_formatting(df: pd.DataFrame, instruction: str, output_path: str) -> str:
    """Applies conditional formatting (highlight cells, color scale) based
    on a natural-language rule, e.g. 'highlight Sales below 0 in red'."""
    schema = f"Columns: {list(df.columns)}"
    system_prompt = (
        "You are an Excel conditional-formatting expert. Reply STRICTLY in "
        'JSON (JSON only): {"column": "Sales", "rule_type": "cell_is", '
        '"operator": "lessThan", "value": 0, "color": "FFC7CE"} '
        'rule_type can be "cell_is" (needs operator: lessThan/greaterThan/equal, '
        'and value) or "color_scale" (no value needed). color is a 6-char '
        "hex code without the #."
    )
    raw = call_llm(system_prompt, f"{schema}\n\nInstruction: {instruction}", max_tokens=300)
    raw = raw.replace("```json", "").replace("```", "").strip()

    import json
    try:
        spec = json.loads(raw)
    except Exception as e:
        return f"Could not parse formatting rule:\n{raw}\n\nError: {e}"

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(list(df.columns))
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")
    for _, row in df.iterrows():
        ws.append(list(row))

    col = spec["column"]
    col_idx = df.columns.get_loc(col) + 1
    col_letter = get_column_letter(col_idx)
    n_rows = len(df) + 1
    cell_range = f"{col_letter}2:{col_letter}{n_rows}"

    if spec["rule_type"] == "color_scale":
        rule = ColorScaleRule(
            start_type="min", start_color="F8696B",
            end_type="max", end_color="63BE7B"
        )
        ws.conditional_formatting.add(cell_range, rule)
    else:
        fill = PatternFill(start_color=spec["color"], end_color=spec["color"], fill_type="solid")
        rule = CellIsRule(operator=spec["operator"], formula=[str(spec["value"])], fill=fill)
        ws.conditional_formatting.add(cell_range, rule)

    wb.save(output_path)
    return f"Conditional formatting applied on '{col}', file saved: {output_path}"


def add_data_validation(df: pd.DataFrame, instruction: str, output_path: str) -> str:
    """Adds a dropdown (data validation) to a column, e.g. restricting
    Status to Active/Inactive/Pending."""
    schema = f"Columns: {list(df.columns)}\nSample:\n{df.head(3).to_string()}"
    system_prompt = (
        "You are an Excel data-validation expert. Reply STRICTLY in JSON "
        '(JSON only): {"column": "Status", "options": ["Active","Inactive","Pending"]}'
    )
    raw = call_llm(system_prompt, f"{schema}\n\nInstruction: {instruction}", max_tokens=300)
    raw = raw.replace("```json", "").replace("```", "").strip()

    import json
    try:
        spec = json.loads(raw)
    except Exception as e:
        return f"Could not parse validation spec:\n{raw}\n\nError: {e}"

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(list(df.columns))
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")
    for _, row in df.iterrows():
        ws.append(list(row))

    col = spec["column"]
    col_idx = df.columns.get_loc(col) + 1
    col_letter = get_column_letter(col_idx)
    n_rows = len(df) + 1
    options_str = ",".join(spec["options"])

    dv = DataValidation(type="list", formula1=f'"{options_str}"', allow_blank=True)
    ws.add_data_validation(dv)
    dv.add(f"{col_letter}2:{col_letter}{n_rows}")

    wb.save(output_path)
    return f"Dropdown validation added on '{col}' ({', '.join(spec['options'])}), file saved: {output_path}"


def edit_excel_in_place(source_path: str, df: pd.DataFrame, output_path: str) -> str:
    """Loads the ORIGINAL Excel file (preserving its sheets, styles, other
    tabs) and updates the main data sheet with the current (possibly
    cleaned) DataFrame, then saves as a new file so the original is safe."""
    if not source_path.lower().endswith((".xlsx", ".xls")):
        return "This command only works when the original file is an Excel (.xlsx) file."

    wb = load_workbook(source_path)
    ws = wb.active  # first/main sheet

    # Clear existing data rows (keep other sheets untouched)
    ws.delete_rows(1, ws.max_row)

    ws.append(list(df.columns))
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF", name="Arial")
        cell.fill = PatternFill("solid", start_color="4472C4")
    for _, row in df.iterrows():
        ws.append(list(row))

    wb.save(output_path)
    sheet_names = wb.sheetnames
    return (
        f"Excel file updated (saved as new file to keep original safe): {output_path}\n"
        f"Sheets preserved: {sheet_names}"
    )


# ----------------------------------------------------------------------
# 6. POWER BI EXPORT (clean, typed data ready for import)
# ----------------------------------------------------------------------
def export_for_powerbi(df: pd.DataFrame, output_path: str) -> str:
    export_df = df.copy()
    for col in export_df.select_dtypes(include="object").columns:
        parsed = pd.to_datetime(export_df[col], errors="coerce")
        if parsed.notna().sum() / max(len(export_df), 1) > 0.9:
            export_df[col] = parsed

    export_df.columns = [str(c).strip().replace(" ", "_") for c in export_df.columns]
    export_df.to_excel(output_path, index=False, sheet_name="Data")
    return f"Power BI-ready file saved: {output_path} (clean headers, typed columns, single flat table)"


# ----------------------------------------------------------------------
# 7. POWER BI: DAX MEASURES + DASHBOARD SETUP GUIDE
# ----------------------------------------------------------------------
def generate_powerbi_guide(df: pd.DataFrame, output_path: str) -> str:
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    cat_cols = df.select_dtypes(include="object").columns.tolist()
    schema = f"Columns: {list(df.columns)}\nNumeric: {numeric_cols}\nCategorical: {cat_cols}\nSample:\n{df.head(3).to_string()}"

    system_prompt = (
        "You are a Power BI consultant. Based on the data schema, provide:\n"
        "1. 5-8 useful DAX measures (copy-paste ready code blocks, each with a short explanation)\n"
        "2. A step-by-step guide on how to build a dashboard from this data in Power BI Desktop — "
        "which visuals (card, bar, line, slicer, matrix) to place where, which column to drag where\n"
        "3. A layout suggestion (e.g. top row KPIs, middle charts, bottom filters)\n"
        "Respond in Markdown format with clear headers, in English "
        "(this guide is saved to a file, so it should stay in professional English)."
        + NO_FILLER_INSTRUCTION
    )

    guide = call_llm(system_prompt, f"Data schema:\n{schema}", max_tokens=1500)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(guide)
    return f"Power BI guide saved: {output_path}\n\n{guide}"


# ----------------------------------------------------------------------
# 9. INTENT DETECTION + GENERAL Q&A (lets the chat understand plain
#    instructions without needing prefixes like 'clean:' / 'sql:')
# ----------------------------------------------------------------------
VALID_INTENTS = {
    "report", "clean", "ask", "sql", "save", "sql_export", "dashboard",
    "formula", "format", "validate", "edit_excel", "powerbi_export",
    "powerbi_guide", "excel_question", "sql_question", "powerbi_question",
    "general_chat",
}


def detect_intent(message: str, has_active_file: bool) -> str:
    """Classifies a free-text chat message into one pipeline action, so the
    user can just type an instruction in plain language instead of using
    prefixes. Falls back to 'general_chat' on any classification failure."""
    system_prompt = (
        "Classify the user's message into exactly ONE of these intents:\n"
        "- report: wants a data quality / EDA report on the uploaded file\n"
        "- clean: wants the data cleaned/modified based on a described rule\n"
        "- ask: wants a specific answer/calculation from the data itself\n"
        "- sql: wants to query the data using SQL / a data question suited to SQL\n"
        "- save: wants the current (cleaned) data saved as CSV/Excel\n"
        "- sql_export: wants the data exported as a standalone SQLite .db file\n"
        "- dashboard: wants an Excel dashboard with charts built\n"
        "- formula: wants a new Excel column driven by a formula\n"
        "- format: wants conditional formatting applied in Excel\n"
        "- validate: wants a dropdown / data validation added in Excel\n"
        "- edit_excel: wants the original Excel file updated in place\n"
        "- powerbi_export: wants a Power BI-ready data export\n"
        "- powerbi_guide: wants DAX measures / a Power BI dashboard-building guide\n"
        "- excel_question: a general Excel question NOT about the uploaded file's data "
        "(e.g. 'how does VLOOKUP work')\n"
        "- sql_question: a general SQL question NOT about the uploaded file's data\n"
        "- powerbi_question: a general Power BI / DAX question NOT about the uploaded file's data\n"
        "- general_chat: anything else, small talk, or questions about the agent itself\n"
        "Reply with ONLY the intent word, nothing else - no punctuation, no explanation."
    )
    user_content = f"An uploaded file is currently loaded: {has_active_file}\nMessage: {message}"
    try:
        raw = call_llm(system_prompt, user_content, max_tokens=15)
        intent = raw.strip().lower().split()[0].strip(".,:;") if raw.strip() else "general_chat"
    except Exception:
        return "general_chat"
    return intent if intent in VALID_INTENTS else "general_chat"


def answer_excel_question(question: str) -> str:
    system_prompt = (
        "You are an Excel expert. Answer precisely and technically, with exact "
        "formula syntax where relevant (e.g. =VLOOKUP(...), =XLOOKUP(...), "
        "=SUMIFS(...)). Respond in English." + NO_FILLER_INSTRUCTION
    )
    return call_llm(system_prompt, question, max_tokens=500)


def answer_sql_question(question: str) -> str:
    system_prompt = (
        "You are a SQL expert. Answer precisely and technically, with SQL "
        "code examples where relevant. Respond in English." + NO_FILLER_INSTRUCTION
    )
    return call_llm(system_prompt, question, max_tokens=500)


def answer_powerbi_question(question: str) -> str:
    system_prompt = (
        "You are a Power BI and DAX expert. Answer precisely and technically, "
        "with DAX or Power Query M code where relevant. Respond in English."
        + NO_FILLER_INSTRUCTION
    )
    return call_llm(system_prompt, question, max_tokens=500)


def chat_reply(message: str) -> str:
    """General conversational fallback - the ONLY place Hinglish + a
    friendlier tone is appropriate, since this is plain chat, not a
    generated artifact."""
    system_prompt = (
        "You are a data analyst assistant. Answer the user's message helpfully."
        + HINGLISH_INSTRUCTION + NO_FILLER_INSTRUCTION
    )
    return call_llm(system_prompt, message, max_tokens=400)


# ----------------------------------------------------------------------
# MAIN CLI LOOP
# ----------------------------------------------------------------------
def main():
    if len(sys.argv) < 2:
        print("Usage: python agent.py <file.csv|xlsx|pdf|docx|pptx>")
        return

    filepath = sys.argv[1]
    print(f"Loading {filepath}...")
    kind, content = load_any_file(filepath)

    if kind == "text":
        print(f"Document loaded ({len(content)} characters)")
        print("\n--- Data Analyst AI Agent ready (document mode) ---")
        print("Commands: '<any question>' | 'exit'\n")
        while True:
            user_input = input("You: ").strip()
            if user_input.lower() in ("exit", "quit"):
                break
            print("\n" + ask_about_text(content, user_input) + "\n")
        return

    # kind == "table"
    df = content
    print(f"Loaded: {df.shape[0]} rows, {df.shape[1]} columns")

    df, log = auto_clean(df)
    print("\nCleaning log:")
    for entry in log:
        print(f"  - {entry}")

    print("\n--- Data Analyst AI Agent ready ---")
    print("Commands: 'report' | 'sql: <question>' | 'clean: <instruction>' | 'save' | 'sql-export' | 'dashboard' | 'powerbi' | 'powerbi-guide' | 'formula: <instruction>' | 'format: <instruction>' | 'validate: <instruction>' | 'edit-excel' | '<any question>' | 'exit'\n")

    while True:
        user_input = input("You: ").strip()
        if user_input.lower() in ("exit", "quit"):
            break
        elif user_input.lower() == "report":
            print("\n" + generate_report(df) + "\n")
        elif user_input.lower().startswith("sql:"):
            print("\n" + run_sql(df, user_input[4:].strip()) + "\n")
        elif user_input.lower().startswith("clean:"):
            df, msg = custom_clean(df, user_input[6:].strip())
            print("\n" + msg + "\n")
        elif user_input.lower() == "save":
            base = os.path.splitext(filepath)[0]
            print("\n" + save_data(df, base) + "\n")
        elif user_input.lower() == "sql-export":
            db_path = os.path.splitext(filepath)[0] + ".db"
            print("\n" + export_sql_db(df, db_path) + "\n")
        elif user_input.lower() == "dashboard":
            out = os.path.splitext(filepath)[0] + "_dashboard.xlsx"
            print("\n" + generate_excel_dashboard(df, out) + "\n")
        elif user_input.lower() == "powerbi":
            out = os.path.splitext(filepath)[0] + "_powerbi.xlsx"
            print("\n" + export_for_powerbi(df, out) + "\n")
        elif user_input.lower() == "powerbi-guide":
            out = os.path.splitext(filepath)[0] + "_powerbi_guide.md"
            print("\n" + generate_powerbi_guide(df, out) + "\n")
        elif user_input.lower().startswith("formula:"):
            out = os.path.splitext(filepath)[0] + "_formula.xlsx"
            print("\n" + add_formula_column(df, user_input[8:].strip(), out) + "\n")
        elif user_input.lower().startswith("format:"):
            out = os.path.splitext(filepath)[0] + "_formatted.xlsx"
            print("\n" + add_conditional_formatting(df, user_input[7:].strip(), out) + "\n")
        elif user_input.lower().startswith("validate:"):
            out = os.path.splitext(filepath)[0] + "_validated.xlsx"
            print("\n" + add_data_validation(df, user_input[9:].strip(), out) + "\n")
        elif user_input.lower() == "edit-excel":
            out = os.path.splitext(filepath)[0] + "_edited.xlsx"
            print("\n" + edit_excel_in_place(filepath, df, out) + "\n")
        else:
            print("\n" + ask_question(df, user_input) + "\n")


if __name__ == "__main__":
    main()